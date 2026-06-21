"""SAYbrand 기말 발표용 PPT 생성 스크립트"""
from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── 색상 ─────────────────────────────────────────────────────
NAVY   = RGBColor(0x0c, 0x14, 0x28)
NAVY2  = RGBColor(0x1d, 0x2a, 0x45)
BLUE   = RGBColor(0x1a, 0x6e, 0xf8)
LBLUE  = RGBColor(0xe8, 0xf0, 0xfb)
RED    = RGBColor(0xdc, 0x26, 0x26)
AMBER  = RGBColor(0xea, 0x58, 0x0c)
YELLOW = RGBColor(0xd9, 0x77, 0x06)
GREEN  = RGBColor(0x16, 0xa3, 0x4a)
WHITE  = RGBColor(0xff, 0xff, 0xff)
GRAY   = RGBColor(0x88, 0x88, 0x88)
LGRAY  = RGBColor(0xf1, 0xf5, 0xf9)
MGRAY  = RGBColor(0xcc, 0xcc, 0xcc)
DGRAY  = RGBColor(0x44, 0x44, 0x44)

TODAY = date.today().isoformat()

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── 헬퍼 ─────────────────────────────────────────────────────
def slide():
    return prs.slides.add_slide(BLANK)


def rect(sl, l, t, w, h, fill):
    shp = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def txt(sl, text, l, t, w, h, size=11, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text or "")
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


def page_header(sl, title, subtitle=""):
    rect(sl, 0, 0, 13.33, 0.82, NAVY)
    rect(sl, 0, 0, 0.06, 7.5, BLUE)
    txt(sl, title, 0.45, 0.1, 11.5, 0.55, size=20, bold=True, color=WHITE)
    if subtitle:
        txt(sl, subtitle, 0.45, 0.58, 11, 0.3, size=10,
            color=RGBColor(0xAA, 0xBB, 0xDD))


def page_footer(sl, page_num):
    rect(sl, 0, 7.32, 13.33, 0.18, NAVY2)
    txt(sl, f"SAYbrand  —  AI 기반 브랜드 위협 모니터링 시스템  |  {TODAY}",
        0.3, 7.34, 10, 0.15, size=7.5, color=RGBColor(0x88, 0x99, 0xbb))
    txt(sl, str(page_num), 12.8, 7.34, 0.45, 0.15,
        size=7.5, color=RGBColor(0x88, 0x99, 0xbb), align=PP_ALIGN.RIGHT)


def section_divider(num, title, subtitle=""):
    sl = slide()
    rect(sl, 0, 0, 13.33, 7.5, NAVY)
    rect(sl, 0, 0, 0.18, 7.5, BLUE)
    txt(sl, num,      0.55, 2.0,  3,  1.0, size=52, bold=True,
        color=RGBColor(0x1a, 0x6e, 0xf8))
    txt(sl, title,    0.55, 3.05, 11, 0.85, size=30, bold=True, color=WHITE)
    if subtitle:
        txt(sl, subtitle, 0.55, 3.98, 10, 0.5,  size=13,
            color=RGBColor(0xAA, 0xBB, 0xDD))
    txt(sl, "SAYbrand", 0.55, 7.1, 4, 0.3, size=8,
        color=RGBColor(0x44, 0x55, 0x77))
    return sl


def row_item(sl, x, y, w, h, label, value, bg=None, label_color=GRAY,
             val_color=NAVY, val_size=10.5, label_size=9):
    if bg:
        rect(sl, x, y, w, h, bg)
    txt(sl, label, x + 0.12, y + 0.06, w * 0.38, h - 0.08,
        size=label_size, color=label_color)
    txt(sl, value, x + w * 0.38 + 0.05, y + 0.06, w * 0.58, h - 0.08,
        size=val_size, bold=True, color=val_color, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 1 — 표지
# ════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, 13.33, 7.5, NAVY)
rect(sl, 0, 0, 13.33, 0.06, BLUE)
rect(sl, 0, 7.44, 13.33, 0.06, BLUE)
rect(sl, 0, 0, 0.18, 7.5, BLUE)

txt(sl, "기말 프로젝트 발표",
    0.55, 0.45, 11, 0.55, size=14,
    color=RGBColor(0xAA, 0xBB, 0xDD))

txt(sl, "SAYbrand",
    0.55, 1.05, 11, 1.5, size=60, bold=True, color=WHITE)

txt(sl, "AI 기반 브랜드 위협 모니터링 시스템",
    0.55, 2.65, 11, 0.65, size=20,
    color=RGBColor(0xAA, 0xBB, 0xDD))

rect(sl, 0.55, 3.5, 11, 0.05, RGBColor(0x1a, 0x6e, 0xf8))

# 공란 — 이름/과목/학교 직접 입력
for i, label in enumerate(["학교", "학과 / 과목명", "팀원"]):
    rect(sl, 0.55, 3.75 + i * 0.75, 11, 0.65, NAVY2)
    txt(sl, label, 0.75, 3.82 + i * 0.75, 2.0, 0.5,
        size=10, color=RGBColor(0xAA, 0xBB, 0xDD))
    txt(sl, "", 2.8, 3.82 + i * 0.75, 8.5, 0.5, size=11, color=WHITE)

txt(sl, TODAY, 0.55, 7.05, 11, 0.35, size=9,
    color=RGBColor(0x44, 0x55, 0x77), align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# Slide 2 — 목차
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "목차")
page_footer(sl, 2)

toc = [
    ("01", "프로젝트 개요",             "무엇을 만들었나"),
    ("02", "문제 정의",                 "왜 만들었나"),
    ("03", "기술 스택 & 개발 환경",     "사용한 도구와 언어"),
    ("04", "시스템 아키텍처",           "전체 설계 구조"),
    ("05", "구현 ① — 3계층 AI 파이프라인", "L1·L2·L3 단계별 구현"),
    ("06", "구현 ② — 리스크 스코어링",  "위협 점수 산출 엔진"),
    ("07", "구현 ③ — 데이터 수집기",   "5개 플랫폼 자동 수집"),
    ("08", "구현 ④ — 카드뉴스 파이프라인", "위협 → 유튜브 쇼츠 자동 생성"),
    ("09", "구현 ⑤ — 보고서 시스템",   "PDF·PPT 자동 생성"),
    ("10", "구현 ⑥ — 조직 관리 & 인증", "RBAC·OAuth·결제"),
    ("11", "완성된 기능 목록",          "구현 현황 총정리"),
    ("12", "기술적 도전과 해결책",      "개발 중 마주한 문제들"),
    ("13", "테스트 결과",               "93 tests passed"),
    ("14", "결론 · 배운 점 · 개선 방향",""),
]

per_col = 7
for i, (num, title, sub) in enumerate(toc):
    col = i // per_col
    row = i % per_col
    x = 0.4 + col * 6.5
    y = 0.92 + row * 0.88
    rect(sl, x, y, 6.18, 0.82, LGRAY)
    rect(sl, x, y, 0.5,  0.82, BLUE if col == 0 else NAVY)
    txt(sl, num,   x + 0.06, y + 0.22, 0.44, 0.4,
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, x + 0.6,  y + 0.06, 5.4,  0.4, size=11, bold=True, color=NAVY)
    if sub:
        txt(sl, sub, x + 0.6, y + 0.46, 5.4, 0.32, size=8.5, color=GRAY)


# ════════════════════════════════════════════════════════════════
# Slide 3 — 프로젝트 개요
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "01  프로젝트 개요", "무엇을 만들었나")
page_footer(sl, 3)

rect(sl, 0.4, 0.95, 12.53, 1.1, NAVY2)
txt(sl, "SAYbrand는 공개 SNS 데이터를 AI로 실시간 분석하여",
    0.6, 1.02, 12.1, 0.48, size=14, bold=True, color=WHITE)
txt(sl, "브랜드를 위협하는 사칭·가짜뉴스·조직적 봇 공격을 자동으로 탐지·대응하는 B2B SaaS입니다.",
    0.6, 1.47, 12.1, 0.48, size=13, color=RGBColor(0xAA, 0xBB, 0xDD))

features = [
    (BLUE,  "수집",   "Naver·YouTube·X\n5개 플랫폼 자동 수집"),
    (GREEN, "분석",   "3계층 AI 파이프라인\nL1→L2→L3 순차 분석"),
    (AMBER, "대응",   "위협 등급별 즉각 알림\nAI 대응 문구 자동 생성"),
    (RED,   "보고",   "일간·주간·월간\nPDF·PPT 자동 생성"),
    (RGBColor(0x7c,0x3a,0xed), "콘텐츠", "위협 데이터 →\n유튜브 쇼츠 자동 생성"),
]
for i, (col, label, desc) in enumerate(features):
    x = 0.4 + i * 2.5
    rect(sl, x, 2.2, 2.35, 2.8, LGRAY)
    rect(sl, x, 2.2, 2.35, 0.5, col)
    txt(sl, label, x + 0.12, 2.25, 2.1, 0.42, size=14, bold=True, color=WHITE)
    txt(sl, desc,  x + 0.12, 2.82, 2.1, 1.9,  size=10.5, color=DGRAY, wrap=True)

rect(sl, 0.4, 5.2, 12.53, 0.42, LGRAY)
txt(sl, "개발 기간:  2026년 상반기   |   배포:  Vercel (API + Frontend)  +  Railway (Celery Worker)   |   버전:  v0.3.1",
    0.6, 5.27, 12.1, 0.3, size=9.5, color=NAVY, align=PP_ALIGN.CENTER)

# 주요 숫자
stats = [("5개", "수집 플랫폼"), ("3계층", "AI 파이프라인"),
         ("900+", "L1 키워드"), ("93개", "테스트 통과"), ("26장", "발표 슬라이드")]
for i, (val, label) in enumerate(stats):
    x = 0.4 + i * 2.5
    rect(sl, x, 5.75, 2.35, 1.42, NAVY2)
    txt(sl, val,   x + 0.12, 5.85, 2.1, 0.72, size=24, bold=True, color=BLUE)
    txt(sl, label, x + 0.12, 6.55, 2.1, 0.5,  size=9.5, color=RGBColor(0xAA,0xBB,0xDD))


# ════════════════════════════════════════════════════════════════
# Slide 4 — 문제 정의
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "02  문제 정의", "왜 이 시스템을 만들었나")
page_footer(sl, 4)

problems = [
    (RED,    "브랜드 사칭",
     "공식 계정을 흉내 낸 가짜 계정이\n소비자를 기만하고 신뢰를 훼손합니다."),
    (AMBER,  "가짜뉴스 · 루머",
     "사실 무근의 악성 루머가 SNS를 타고\n수십만 명에게 순식간에 확산됩니다."),
    (YELLOW, "조직적 봇 공격",
     "경쟁사 또는 악의적 세력이 봇 계정으로\n집단적 부정 댓글·리뷰 공격을 합니다."),
    (BLUE,   "임직원 리스크",
     "임원·직원의 개인 SNS 발언이\n기업 이미지에 치명타를 입힙니다."),
]
for i, (col, title, desc) in enumerate(problems):
    x = 0.4 + i * 3.23
    rect(sl, x, 0.95, 3.08, 4.2, LGRAY)
    rect(sl, x, 0.95, 3.08, 0.08, col)
    rect(sl, x, 0.95, 0.08, 4.2,  col)
    txt(sl, title, x + 0.2, 1.12, 2.78, 0.52, size=13, bold=True, color=NAVY)
    txt(sl, desc,  x + 0.2, 1.72, 2.78, 2.0,  size=10.5, color=DGRAY, wrap=True)

rect(sl, 0.4, 5.35, 12.53, 0.06, BLUE)
txt(sl, "핵심 문제",
    0.4, 5.55, 2.2, 0.45, size=12, bold=True, color=NAVY)
txt(sl, "SNS 데이터는 너무 방대하고 빠르게 생성되어 사람이 24시간 수동으로 모니터링하는 것은 불가능합니다.\n"
       "또한 봇 공격인지 실제 소비자 불만인지 구분하지 못하면 불필요한 법적 대응으로 비용이 낭비됩니다.",
    0.4, 5.55, 12.53, 1.4, size=10.5, color=DGRAY, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 5 — 기술 스택
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "03  기술 스택 & 개발 환경")
page_footer(sl, 5)

stack_cols = [
    ("백엔드", NAVY, [
        ("언어 · 프레임워크", "Python 3.11+  /  FastAPI 0.115"),
        ("ORM · DB",          "SQLAlchemy 2.0 async  /  SQLite(로컬) · PostgreSQL(운영)"),
        ("비동기 워커",       "Celery 5.4  +  Redis"),
        ("인증",              "Google OAuth 2.0  (Authlib)"),
        ("결제",              "Polar  +  Svix 웹훅 서명 검증"),
    ]),
    ("AI · 분석", BLUE, [
        ("L1 — 규칙 필터",   "자체 키워드 DB (900개+, 18개 카테고리)"),
        ("L2 — 텍스트 분석", "Gemini 2.5 Flash Lite  (google-genai ≥ 1.0)"),
        ("L3 — 심층 분석",   "Claude Haiku 4.5  (anthropic SDK)"),
        ("감성 폴백",         "KNU 한국어 감성 사전 (14,854 단어)"),
        ("이미지",            "imagehash  (pHash)  /  Google Vision API"),
    ]),
    ("프론트엔드 · 인프라", GREEN, [
        ("CSS · JS",          "Tailwind CSS CDN  +  Vanilla JS  (빌드 없음)"),
        ("폰트",              "Syne  /  Noto Sans KR  /  JetBrains Mono"),
        ("PWA",               "manifest.json  +  Service Worker"),
        ("배포 — API",        "Vercel  (진입점: app.py)"),
        ("배포 — 워커",       "Railway  (Celery -B worker)"),
    ]),
]

for ci, (group, col, items) in enumerate(stack_cols):
    x = 0.4 + ci * 4.3
    rect(sl, x, 0.95, 4.1, 0.42, col)
    txt(sl, group, x + 0.12, 0.99, 3.86, 0.35, size=12, bold=True, color=WHITE)
    for ri, (label, val) in enumerate(items):
        bg = LGRAY if ri % 2 == 0 else WHITE
        rect(sl, x, 1.37 + ri * 1.12, 4.1, 1.1, bg)
        txt(sl, label, x + 0.12, 1.42 + ri * 1.12, 1.6,  0.38,
            size=8.5, color=GRAY)
        txt(sl, val,   x + 0.12, 1.78 + ri * 1.12, 3.85, 0.58,
            size=9.5, bold=True, color=NAVY, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 6 — 시스템 아키텍처
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "04  시스템 아키텍처", "전체 데이터 흐름")
page_footer(sl, 6)

# 흐름도 — 수평 배치
arch = [
    (BLUE,  "SNS\n플랫폼",   "YouTube\nNaver · X"),
    (NAVY,  "수집기",        "Celery Beat\n30분 주기"),
    (GREEN, "L1 필터",       "규칙 기반\n$0 비용"),
    (BLUE,  "L2 분석",       "Gemini\nFlash Lite"),
    (AMBER, "L3 심층",       "Claude\nHaiku 4.5"),
    (RED,   "DB 저장",       "PostgreSQL\n/ SQLite"),
]
for i, (col, title, sub) in enumerate(arch):
    x = 0.38 + i * 2.12
    rect(sl, x, 0.95, 1.95, 1.55, col)
    txt(sl, title, x + 0.1,  1.02, 1.75, 0.6,  size=12, bold=True, color=WHITE)
    txt(sl, sub,   x + 0.1,  1.6,  1.75, 0.75, size=9.5, color=WHITE)
    if i < 5:
        txt(sl, "→", x + 1.96, 1.55, 0.2, 0.45, size=14, bold=True, color=NAVY)

# 분기 화살표
rect(sl, 0.38, 2.65, 12.57, 0.06, MGRAY)
for xi, label in [(0.5, "대시보드 API"), (4.5, "Slack / 이메일 알림"),
                  (8.5, "PDF · PPT 보고서"), (10.8, "카드뉴스 쇼츠")]:
    txt(sl, "↓", xi + 0.6, 2.72, 0.3, 0.4, size=13, color=BLUE, align=PP_ALIGN.CENTER)
    txt(sl, label, xi, 3.1, 2.1, 0.45, size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# 프론트엔드 박스
rect(sl, 0.38, 3.65, 5.5, 2.78, LGRAY)
rect(sl, 0.38, 3.65, 5.5, 0.4,  NAVY)
txt(sl, "프론트엔드 (Vanilla JS + Tailwind CDN)", 0.52, 3.7, 5.2, 0.32,
    size=10, bold=True, color=WHITE)
pages = ["대시보드  /  위협 목록  /  행동 목록",
         "보고서  /  설정  /  고객센터",
         "랜딩  /  온보딩  /  조직 관리"]
for pi, p in enumerate(pages):
    bg = WHITE if pi % 2 == 0 else LGRAY
    rect(sl, 0.38, 4.05 + pi * 0.75, 5.5, 0.72, bg)
    txt(sl, p, 0.52, 4.12 + pi * 0.75, 5.2, 0.58, size=9.5, color=DGRAY)

# 인프라 박스
rect(sl, 6.2, 3.65, 6.75, 2.78, LGRAY)
rect(sl, 6.2, 3.65, 6.75, 0.4,  NAVY2)
txt(sl, "배포 인프라", 6.35, 3.7, 6.4, 0.32, size=10, bold=True, color=WHITE)
infra = [
    (BLUE,  "Vercel",   "API 라우터 + 정적 프론트엔드 서빙  (진입점: app.py)"),
    (RED,   "Railway",  "Celery 워커 + Beat 스케줄러  (-B -c 2)"),
    (GREEN, "Redis",    "Celery 브로커 + 결과 백엔드"),
    (AMBER, "PostgreSQL","운영 DB  (로컬: SQLite)"),
]
for ii, (col, name, desc) in enumerate(infra):
    bg = WHITE if ii % 2 == 0 else LGRAY
    rect(sl, 6.2, 4.05 + ii * 0.6, 6.75, 0.58, bg)
    rect(sl, 6.2, 4.05 + ii * 0.6, 0.08, 0.58, col)
    txt(sl, name, 6.35, 4.1 + ii * 0.6, 1.5,  0.45, size=9.5, bold=True, color=NAVY)
    txt(sl, desc, 7.9,  4.1 + ii * 0.6, 4.9,  0.45, size=9, color=DGRAY, wrap=True)


# ════════════════════════════════════════════════════════════════
# Section 05
# ════════════════════════════════════════════════════════════════
section_divider("05", "구현 ①\n3계층 AI 파이프라인", "비용을 최소화하면서 정확도를 극대화")


# ════════════════════════════════════════════════════════════════
# Slide 7 — AI 파이프라인
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "구현 ① — 3계층 AI 파이프라인",
            "L1($0) → L2(저비용) → L3(고위협만) 순서로 AI 비용 최소화")
page_footer(sl, 7)

layers = [
    (GREEN, "L1", "규칙 기반 필터", "비용 $0",
     "backend/services/analyzers/l1_filter.py",
     ["900개+ 키워드 DB (18개 카테고리)",
      "CRITICAL_BYPASS: 법적 위협 패턴 즉시 통과",
      "NEGATIVE_FILTERS 20개+ — 오탐(범용어) 방지",
      "score < 0.05 → 탈락, AI 호출 없음",
      "industry_threshold: food=0.067 / general=0.08"]),
    (BLUE,  "L2", "감성·의도 분석", "배치 10건/호출",
     "backend/services/analyzers/l2_text.py",
     ["HyperCLOVA X → Gemini 2.5 Flash Lite → KNU 폴백",
      "12개 마케팅 위기 카테고리 분류",
      "감성(positive/neutral/negative) + 감정 7가지",
      "봇 확률(0.0–1.0) + 조직적 공격 여부",
      "불완전 JSON 배열 자동 복구 로직"]),
    (AMBER, "L3", "심층 대응 분석", "고위협(score≥85)만",
     "backend/services/analyzers/l3_deep.py",
     ["Gemini 2.5 Flash Lite → Claude Haiku 4.5 폴백",
      "brand_damage_type 분류",
      "  (매출타격·채용악영향·파트너십위험 등)",
      "communication_urgency (즉시/당일/48h/모니터링)",
      "SNS대응·보도자료·내부조치 대응 문구 자동 생성"]),
]
for i, (col, lnum, ltitle, cost, fpath, bullets) in enumerate(layers):
    x = 0.38 + i * 4.3
    rect(sl, x, 0.92, 4.1, 6.2, LGRAY)
    rect(sl, x, 0.92, 4.1, 0.52, col)
    txt(sl, lnum,   x + 0.12, 0.96, 0.72, 0.42, size=18, bold=True, color=WHITE)
    txt(sl, ltitle, x + 0.78, 0.96, 3.1,  0.42, size=12, bold=True, color=WHITE)
    rect(sl, x, 1.44, 4.1, 0.32, RGBColor(0xdd,0xdd,0xdd))
    txt(sl, cost,   x + 0.12, 1.48, 2.0, 0.26, size=8.5, bold=True, color=DGRAY)
    txt(sl, fpath,  x + 0.12, 1.76, 3.85, 0.3,  size=7.5, italic=True, color=GRAY)
    for bi, b in enumerate(bullets):
        bg = WHITE if bi % 2 == 0 else LGRAY
        rect(sl, x, 2.1 + bi * 0.92, 4.1, 0.9, bg)
        txt(sl, ("• " if not b.startswith(" ") else "") + b,
            x + 0.12, 2.15 + bi * 0.92, 3.85, 0.8,
            size=9.5, color=DGRAY, wrap=True)
    if i < 2:
        txt(sl, "→", x + 4.11, 3.8, 0.22, 0.45, size=14, bold=True, color=col)


# ════════════════════════════════════════════════════════════════
# Section 06
# ════════════════════════════════════════════════════════════════
section_divider("06", "구현 ②\n리스크 스코어링 엔진", "위협 점수 0–100 자동 산출")


# ════════════════════════════════════════════════════════════════
# Slide 8 — 리스크 스코어링
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "구현 ② — 리스크 스코어링 엔진",
            "backend/services/risk_scorer.py  — 설계에서 확정된 가중치 테이블")
page_footer(sl, 8)

rect(sl, 0.38, 0.92, 12.57, 1.05, NAVY2)
txt(sl, "base = SEVERITY_WEIGHTS[severity] × MODULE_WEIGHTS[module] × PLATFORM_WEIGHTS[platform] × confidence × 100",
    0.58, 0.98, 12.2, 0.45, size=10.5, bold=True, color=WHITE)
txt(sl, "base  ×=  industry_multiplier     if is_organized: base = min(base × 1.3, 100)     final = base × (recency_weight + velocity_bonus)",
    0.58, 1.42, 12.2, 0.42, size=9.5, color=RGBColor(0xAA,0xBB,0xDD))

tables = [
    ("심각도 (SEVERITY_WEIGHTS)", [
        ("critical", "1.0",  RED),
        ("high",     "0.7",  AMBER),
        ("medium",   "0.4",  YELLOW),
        ("low",      "0.15", GREEN),
    ]),
    ("모듈 (MODULE_WEIGHTS)", [
        ("A — 브랜드 사칭",     "1.0",  RED),
        ("B — 루머·가짜뉴스",  "0.85", AMBER),
        ("C — 임직원 리스크",  "0.7",  YELLOW),
    ]),
    ("플랫폼 (PLATFORM_WEIGHTS)", [
        ("Instagram",  "1.0", NAVY),
        ("YouTube",    "0.9", NAVY),
        ("TikTok",     "0.85",NAVY),
        ("X (Twitter)","0.8", NAVY),
        ("Naver",      "0.7", NAVY),
    ]),
]
for ti, (title, rows) in enumerate(tables):
    x = 0.38 + ti * 4.3
    rect(sl, x, 2.1, 4.1, 0.42, NAVY)
    txt(sl, title, x + 0.1, 2.14, 3.9, 0.35, size=10, bold=True, color=WHITE)
    for ri, (lbl, val, col) in enumerate(rows):
        bg = LGRAY if ri % 2 == 0 else WHITE
        rect(sl, x, 2.52 + ri * 0.52, 4.1, 0.52, bg)
        txt(sl, lbl, x + 0.1, 2.57 + ri * 0.52, 3.0, 0.42, size=9.5, color=DGRAY)
        txt(sl, val, x + 3.3, 2.57 + ri * 0.52, 0.7, 0.42, size=10, bold=True, color=col, align=PP_ALIGN.RIGHT)

rect(sl, 0.38, 5.4, 12.57, 0.42, LGRAY)
txt(sl, "최신성 가중치:  < 1h → 1.0  |  1–6h → 0.9  |  6–24h → 0.75  |  1–3일 → 0.5  |  3–7일 → 0.3  |  7일+ → 0.1",
    0.55, 5.46, 12.2, 0.32, size=9.5, color=NAVY, align=PP_ALIGN.CENTER)

rect(sl, 0.38, 5.95, 12.57, 0.42, NAVY2)
txt(sl, "임계값:  80–100 CRITICAL (즉각 Slack)   60–79 HIGH (당일)   35–59 MEDIUM (모니터링)   0–34 LOW (리포트)",
    0.55, 6.01, 12.2, 0.32, size=9.5, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "조직적 공격 탐지: text_uniformity·account_cluster·temporal_cluster 등 6개 지표 가중합 → attack_score ≥ 0.7 → organized_attack",
    0.38, 6.5, 12.57, 0.65, size=9.5, color=DGRAY, wrap=True)


# ════════════════════════════════════════════════════════════════
# Section 07
# ════════════════════════════════════════════════════════════════
section_divider("07", "구현 ③\n데이터 수집기", "5개 플랫폼 Celery 자동 수집")


# ════════════════════════════════════════════════════════════════
# Slide 9 — 수집기
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "구현 ③ — 데이터 수집기",
            "backend/services/collectors/  — Celery Beat 30분 주기 병렬 수집")
page_footer(sl, 9)

rect(sl, 0.38, 0.92, 12.57, 0.42, NAVY)
for hdr, xp in [("상태", 0.55), ("플랫폼", 1.1),
                 ("수집 대상", 2.85), ("구현 방법", 5.5), ("비고", 10.0)]:
    txt(sl, hdr, xp, 0.97, 2.2, 0.32, size=9.5, bold=True, color=WHITE)

rows = [
    (GREEN,  "✅ 실동작", "Naver",       "블로그·카페·뉴스·이미지",
     "NAVER_CLIENT_ID/SECRET → Search API",       "검증 완료"),
    (GREEN,  "✅ API 키", "YouTube",     "영상 댓글·메타데이터",
     "YOUTUBE_API_KEY → Data API v3",             "remove_pii 적용"),
    (GREEN,  "✅ API 키", "X (Twitter)", "트윗·대화",
     "X_BEARER_TOKEN → v2 API",                   "remove_pii 적용"),
    (YELLOW, "🟡 Mock",  "한국 커뮤니티","에펨코리아·더쿠·클리앙",
     "크롤링 (robots.txt 자동 준수)",              "Mock 처리"),
    (GRAY,   "❌ v1.1",  "Instagram",   "게시글·댓글·스토리",
     "Meta API 접근 제한",                         "v1.1 구현 예정"),
    (GRAY,   "❌ v1.1",  "TikTok",      "영상·댓글",
     "TikTok API 접근 제한",                       "v1.1 구현 예정"),
]
for ri, (col, status, plat, target, method, note) in enumerate(rows):
    bg = LGRAY if ri % 2 == 0 else WHITE
    y = 1.34 + ri * 0.77
    rect(sl, 0.38, y, 12.57, 0.75, bg)
    rect(sl, 0.38, y, 0.1,   0.75, col)
    txt(sl, status, 0.52,  y + 0.18, 0.62, 0.42, size=9,  color=col)
    txt(sl, plat,   1.1,   y + 0.18, 1.75, 0.42, size=10, bold=True, color=NAVY)
    txt(sl, target, 2.85,  y + 0.18, 2.65, 0.42, size=9.5, color=DGRAY)
    txt(sl, method, 5.5,   y + 0.18, 4.5,  0.42, size=9, color=DGRAY)
    txt(sl, note,   10.0,  y + 0.18, 2.8,  0.42, size=9.5, bold=(col==GREEN), color=col)

rect(sl, 0.38, 5.98, 12.57, 0.82, LGRAY)
txt(sl, "공통 컴플라이언스 처리",
    0.55, 6.03, 4.0, 0.32, size=10, bold=True, color=NAVY)
txt(sl, "robots.txt 자동 체크 → 수집 전 허용 여부 확인  |  전화번호·이메일·주민번호 정규식 PII 마스킹\n"
       "요청 간 최소 2초 지연 (Rate Limiter)  |  뉴스 도메인 자동 분류 (is_news_domain())",
    0.55, 6.35, 12.2, 0.42, size=9, color=DGRAY, wrap=True)


# ════════════════════════════════════════════════════════════════
# Section 08
# ════════════════════════════════════════════════════════════════
section_divider("08", "구현 ④\n카드뉴스 파이프라인", "위협 데이터 → 유튜브 쇼츠 영상 완전 자동화")


# ════════════════════════════════════════════════════════════════
# Slide 10 — 카드뉴스 파이프라인
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "구현 ④ — 카드뉴스 자동 생성 파이프라인",
            "card-news-pipeline/  — SAYbrand 위협 DB → 유튜브 쇼츠 영상 완전 자동화")
page_footer(sl, 10)

steps = [
    (NAVY,  "① DB 로드",
     "db_source.py",
     "최근 14일 위협\ncritical/high/medium\n최신순 DESC LIMIT 100\nMock AI 분석 필터링"),
    (BLUE,  "② 소재 선택",
     "selector.py",
     "오늘 우선 탐색\n→ 1·2·3일 전\n→ 전체 최고점\nused_ids 중복 제거"),
    (GREEN, "③ 스크립팅",
     "llm_scripter.py",
     "Claude Haiku 4.5\nheadline ≤ 20자\nbody ≤ 150자\n태그 5개 JSON"),
    (AMBER, "④ 렌더링",
     "renderer.py",
     "Playwright Chromium\n1080×1920 세로형\nPixazo 히어로 이미지\nPNG 슬라이드 저장"),
    (RED,   "⑤ 영상 조립",
     "video.py",
     "FFmpeg MP4 변환\nassets/bgm/*.mp3\n배경음악 믹싱\n출력: source_id.mp4"),
    (DGRAY, "⑥ 검수·업로드",
     "orchestrator.py",
     "Discord Webhook\n검수 요청 전송\n승인 후 YouTube\n비공개 Shorts 업로드"),
]
for i, (col, title, fpath, desc) in enumerate(steps):
    x = 0.38 + i * 2.16
    rect(sl, x, 0.92, 2.0,  5.78, LGRAY)
    rect(sl, x, 0.92, 2.0,  0.48, col)
    txt(sl, title, x + 0.1, 0.96, 1.8, 0.38, size=11, bold=True, color=WHITE)
    txt(sl, fpath, x + 0.1, 1.44, 1.8, 0.3,  size=7.5, italic=True, color=GRAY)
    txt(sl, desc,  x + 0.1, 1.82, 1.8, 4.5,  size=9.5, color=DGRAY, wrap=True)
    if i < 5:
        txt(sl, "→", x + 2.01, 3.6, 0.18, 0.42, size=12, bold=True, color=BLUE)

rect(sl, 0.38, 6.82, 12.57, 0.35, NAVY2)
txt(sl, "폴백:  ANTHROPIC_API_KEY 없을 때 scripter.py 규칙 기반 템플릿 자동 대체  |  "
       "DATABASE_URL 없을 때 mock_data.py 샘플 사용",
    0.55, 6.88, 12.2, 0.28, size=8.5, color=WHITE)


# ════════════════════════════════════════════════════════════════
# Section 09
# ════════════════════════════════════════════════════════════════
section_divider("09", "구현 ⑤\n보고서 시스템", "일간·주간·월간 PDF / PPT 자동 생성")


# ════════════════════════════════════════════════════════════════
# Slide 11 — 보고서 시스템
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "구현 ⑤ — 보고서 시스템",
            "backend/services/report_generator.py  +  pptx_generator.py")
page_footer(sl, 11)

# PDF
rect(sl, 0.38, 0.92, 6.0, 0.45, RED)
txt(sl, "PDF 보고서  —  GET /api/reports/{period}/pdf",
    0.52, 0.97, 5.72, 0.36, size=11, bold=True, color=WHITE)

pdf_items = [
    ("라이브러리",  "ReportLab ≥ 4.0"),
    ("용지",        "A4, 8섹션 구성"),
    ("한국어 폰트", "NanumGothic TTF 번들 (서버 폰트 없을 때 자동 탐색)"),
    ("텍스트 처리", "wordWrap='CJK' — 한국어 셀 오버플로 방지"),
    ("레이아웃",    "테이블 너비 180mm (A4 210mm − 좌우 여백 각 15mm)"),
    ("자동 삽입",   "헤더·푸터 (표지 제외) + 페이지 번호"),
    ("섹션 구성",   "표지·핵심요약·위협현황·감성·조직공격·Top10·조치·권고"),
]
for ri, (lbl, val) in enumerate(pdf_items):
    bg = LGRAY if ri % 2 == 0 else WHITE
    rect(sl, 0.38, 1.37 + ri * 0.72, 6.0, 0.7, bg)
    txt(sl, lbl, 0.52, 1.42 + ri * 0.72, 1.8, 0.58, size=8.5, color=GRAY)
    txt(sl, val, 2.35, 1.42 + ri * 0.72, 3.9, 0.58, size=9.5, bold=True, color=NAVY, wrap=True)

# PPT
rect(sl, 6.73, 0.92, 6.22, 0.45, BLUE)
txt(sl, "PPT 보고서  —  GET /api/reports/{period}/pptx",
    6.87, 0.97, 5.94, 0.36, size=11, bold=True, color=WHITE)

ppt_items = [
    ("라이브러리",  "python-pptx ≥ 0.6.21"),
    ("규격",        "16:9 (13.33 × 7.5 인치), 6슬라이드"),
    ("슬라이드 1",  "Cover — NAVY 배경, KPI 4개 박스"),
    ("슬라이드 2",  "핵심 요약 — KPI 6개 그리드 + 브랜드 점수"),
    ("슬라이드 3",  "위협 현황 — 심각도·플랫폼 인라인 바 차트"),
    ("슬라이드 4",  "감성 분석 — 감성 분포 + 감정 분류"),
    ("슬라이드 5·6","미해결 Top5 테이블 + 대응 권고사항 카드"),
]
for ri, (lbl, val) in enumerate(ppt_items):
    bg = LGRAY if ri % 2 == 0 else WHITE
    rect(sl, 6.73, 1.37 + ri * 0.72, 6.22, 0.7, bg)
    txt(sl, lbl, 6.87, 1.42 + ri * 0.72, 1.8,  0.58, size=8.5, color=GRAY)
    txt(sl, val, 8.7,  1.42 + ri * 0.72, 4.12, 0.58, size=9.5, bold=True, color=NAVY, wrap=True)

rect(sl, 0.38, 6.47, 12.57, 0.72, LGRAY)
txt(sl, "공통:  period = daily / weekly / monthly  |  is_mock 필드로 실데이터 여부 구분  |  "
       "브랜드 점수·KPI·감성·감정·조직공격·봇 수 등 통합 JSON 리포트 먼저 생성 후 PDF/PPT 렌더링",
    0.55, 6.53, 12.2, 0.6, size=9, color=DGRAY, wrap=True)


# ════════════════════════════════════════════════════════════════
# Section 10
# ════════════════════════════════════════════════════════════════
section_divider("10", "구현 ⑥\n조직 관리 & 인증", "RBAC · Google OAuth · Polar 결제")


# ════════════════════════════════════════════════════════════════
# Slide 12 — 조직 관리 & 인증
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "구현 ⑥ — 조직 관리 & 인증 & 결제",
            "backend/routers/  auth · orgs · billing")
page_footer(sl, 12)

# RBAC
rect(sl, 0.38, 0.92, 3.9, 0.42, NAVY)
txt(sl, "멤버 역할 (RBAC)", 0.52, 0.97, 3.62, 0.34, size=10.5, bold=True, color=WHITE)
roles = [("owner","전체 관리 (삭제 포함)",RED),
         ("admin","멤버 관리·설정 변경",AMBER),
         ("member","스캔·위협 처리",BLUE),
         ("viewer","읽기 전용",GRAY)]
for ri, (role, perm, col) in enumerate(roles):
    bg = LGRAY if ri % 2 == 0 else WHITE
    rect(sl, 0.38, 1.34 + ri * 0.6, 3.9, 0.58, bg)
    rect(sl, 0.38, 1.34 + ri * 0.6, 0.08, 0.58, col)
    txt(sl, role, 0.52, 1.39 + ri * 0.6, 1.2, 0.46, size=9.5, bold=True, color=col)
    txt(sl, perm, 1.75, 1.39 + ri * 0.6, 2.4, 0.46, size=9,   color=DGRAY)

# 가입 플로우
rect(sl, 4.5, 0.92, 4.2, 0.42, NAVY)
txt(sl, "가입 플로우", 4.65, 0.97, 3.9, 0.34, size=10.5, bold=True, color=WHITE)
flows = [
    (GREEN, "초대 코드",
     "관리자 → 코드 생성(역할·만료·횟수)\n→ URL 공유 → 입력 → 즉시 active"),
    (BLUE,  "승인 요청",
     "사용자 → 참여 신청(pending)\n→ 관리자 승인(active) / 거절(DB 삭제)"),
]
for fi, (col, ftitle, fdesc) in enumerate(flows):
    rect(sl, 4.5, 1.34 + fi * 1.28, 4.2, 0.42, col)
    txt(sl, ftitle, 4.65, 1.38 + fi * 1.28, 3.9, 0.34, size=10, bold=True, color=WHITE)
    rect(sl, 4.5, 1.76 + fi * 1.28, 4.2, 0.82, LGRAY)
    txt(sl, fdesc,  4.65, 1.8  + fi * 1.28, 3.9, 0.75, size=9.5, color=DGRAY, wrap=True)

# 인증 & 결제
rect(sl, 8.92, 0.92, 4.03, 0.42, NAVY)
txt(sl, "인증 & 결제", 9.06, 0.97, 3.75, 0.34, size=10.5, bold=True, color=WHITE)
auth_items = [
    ("Google OAuth 2.0", "Authlib — /auth/google → 콜백 → JWT 발급"),
    ("Polar 결제",        "/billing/checkout/{plan} → 체크아웃 링크 리다이렉트"),
    ("Svix 웹훅",         "whsec_ prefix 제거 후 base64 decode 서명 검증"),
    ("구독 동기화",       "POST /billing/sync — 이메일 기준 수동 동기화"),
    ("티어 제한",         "free=조직1  starter=3  pro=5  enterprise=무제한"),
]
for ri, (lbl, val) in enumerate(auth_items):
    bg = LGRAY if ri % 2 == 0 else WHITE
    rect(sl, 8.92, 1.34 + ri * 0.6, 4.03, 0.58, bg)
    txt(sl, lbl, 9.06, 1.38 + ri * 0.6, 1.7,  0.46, size=8.5, color=GRAY)
    txt(sl, val, 10.8, 1.38 + ri * 0.6, 2.05, 0.46, size=8.5, color=NAVY, wrap=True)

# 온보딩
rect(sl, 0.38, 3.8, 12.57, 0.42, NAVY2)
txt(sl, "3단계 온보딩 (/onboarding)  — CustomerProfile 없으면 자동 리다이렉트",
    0.55, 3.85, 12.2, 0.32, size=10, bold=True, color=WHITE)
ob_steps = [
    ("Step 1", "경로 선택", "새 조직 만들기 / 기존 조직 참여"),
    ("Step 2", "조직 설정", "조직명·도메인 / 초대 코드 입력"),
    ("Step 3", "브랜드 등록", "브랜드명·별칭·업종·SNS계정·임직원·키워드"),
]
for si, (num, title, desc) in enumerate(ob_steps):
    x = 0.38 + si * 4.26
    rect(sl, x, 4.22, 4.1, 2.85, LGRAY)
    rect(sl, x, 4.22, 4.1, 0.42, BLUE)
    txt(sl, num,   x + 0.12, 4.26, 0.9,  0.35, size=10, bold=True, color=WHITE)
    txt(sl, title, x + 0.95, 4.26, 3.0,  0.35, size=10, bold=True, color=WHITE)
    txt(sl, desc,  x + 0.12, 4.74, 3.85, 2.2,  size=9.5, color=DGRAY, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 13 — 완성된 기능 목록
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "11  완성된 기능 목록", f"구현 현황 총정리  —  기준일 {TODAY}")
page_footer(sl, 13)

items = [
    ("✅", GREEN,  "코어 AI 파이프라인 (L1→L2→L3)",    "실동작 검증"),
    ("✅", GREEN,  "리스크 스코어링 엔진",               "가중치 테이블 확정"),
    ("✅", GREEN,  "Naver 수집기",                       "실 API 검증 완료"),
    ("✅", GREEN,  "YouTube · X 수집기",                 "API 키 있을 때 실동작"),
    ("✅", GREEN,  "Gemini L2 분석",                     "유료 전환 완료"),
    ("✅", GREEN,  "Claude Haiku L3 분석",               "API 키 있을 때 실동작"),
    ("✅", GREEN,  "KNU 감성 사전 폴백",                 "오프라인 항상 동작"),
    ("✅", GREEN,  "Google OAuth 2.0 인증",              "실동작"),
    ("✅", GREEN,  "Polar 결제 + Svix 웹훅",             "서명 검증 실동작"),
    ("✅", GREEN,  "조직 관리 (RBAC·초대코드·승인)",    "실동작"),
    ("✅", GREEN,  "3단계 온보딩 플로우",                "실동작"),
    ("✅", GREEN,  "PDF / PPT 보고서 자동 생성",         "일간·주간·월간"),
    ("✅", GREEN,  "카드뉴스 파이프라인",                "93 tests passed"),
    ("✅", GREEN,  "Slack 알림",                         "위협 등급별 분기"),
    ("✅", GREEN,  "고객센터 게시판",                    "support_admin_emails 권한"),
    ("✅", GREEN,  "PWA (manifest + Service Worker)",    "모바일 설치 가능"),
    ("🟡", AMBER,  "Celery 비동기 워커",                 "Redis 환경 필요"),
    ("🟡", AMBER,  "이미지 분석 (pHash)",                "엔진 구현, 파이프라인 미통합"),
    ("❌", GRAY,   "Instagram 수집기",                   "Meta API 제한 — v1.1"),
    ("❌", GRAY,   "TikTok 수집기",                     "API 제한 — v1.1"),
]

per_col = 10
for i, (icon, col, feature, note) in enumerate(items):
    ci = i // per_col
    ri = i % per_col
    x = 0.38 + ci * 6.48
    y = 0.92 + ri * 0.62
    bg = LGRAY if ri % 2 == 0 else WHITE
    rect(sl, x, y, 6.3, 0.6, bg)
    txt(sl, icon,    x + 0.08, y + 0.1,  0.38, 0.42, size=12)
    txt(sl, feature, x + 0.52, y + 0.1,  3.95, 0.42, size=9.5, color=NAVY)
    txt(sl, note,    x + 4.55, y + 0.1,  1.65, 0.42, size=8.5, color=col, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# Slide 14 — 기술적 도전과 해결책
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "12  기술적 도전과 해결책", "개발 과정에서 마주한 주요 문제들")
page_footer(sl, 14)

challenges = [
    (RED,   "Vercel lifespan 미작동",
     "문제: FastAPI lifespan 이벤트가 Vercel 서버리스 환경에서 실행되지 않아 DB 자동 마이그레이션 불가",
     "해결: Railway PostgreSQL에서 직접 SQL 실행 / app.py를 Vercel 진입점으로 분리"),
    (AMBER, "datetime 타입 불일치",
     "문제: datetime.now(timezone.utc) 사용 시 asyncpg DataError 발생 — aware datetime을 거부",
     "해결: DB 전체를 naive UTC(datetime.utcnow())로 통일"),
    (BLUE,  "Polar Svix 웹훅 서명 검증 실패",
     "문제: 헤더 형식이 'v1,<base64>' (쉼표, hex 아님) — whsec_ 접두어 포함 시 검증 실패",
     "해결: whsec_ prefix 제거 후 base64 decode, 헤더 파싱 방식 수정"),
    (GREEN, "L2 배치 응답 불완전 JSON",
     "문제: Gemini가 배치 10건 처리 시 토큰 한계로 JSON 배열 중간 절단",
     "해결: max_output_tokens 증가 + 불완전 배열 자동 복구 로직 구현"),
    (NAVY,  "dismissed_urls 쿼리 버그",
     "문제: Python bool과 SQLAlchemy BinaryExpression 혼용으로 스캔 결과 0건 반환",
     "해결: SQLAlchemy 표현식 방식으로 통일 (commit d0bc54b)"),
    (RGBColor(0x7c,0x3a,0xed), "Gemini 모델명 deprecated",
     "문제: gemini-2.0-flash-lite 호출 시 404 오류 — 이미 deprecated된 모델명 사용",
     "해결: gemini-2.5-flash-lite로 교체 + google-genai ≥ 1.0.0 SDK로 마이그레이션"),
]

for i, (col, title, prob, sol) in enumerate(challenges):
    row = i // 2
    ci  = i % 2
    x = 0.38 + ci * 6.48
    y = 0.92 + row * 2.12
    rect(sl, x, y,        6.3, 2.08, LGRAY)
    rect(sl, x, y,        6.3, 0.38, col)
    rect(sl, x, y,        0.1, 2.08, col)
    txt(sl, title, x + 0.18, y + 0.05, 6.0, 0.3, size=11, bold=True, color=WHITE)
    txt(sl, "문제: " + prob.replace("문제: ",""), x + 0.18, y + 0.46, 6.0, 0.75,
        size=9, color=DGRAY, wrap=True)
    txt(sl, "해결: " + sol.replace("해결: ",""), x + 0.18, y + 1.22, 6.0, 0.78,
        size=9, bold=True, color=NAVY, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 15 — 테스트 결과
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "13  테스트 결과", "pytest — 카드뉴스 파이프라인 93 tests passed")
page_footer(sl, 15)

rect(sl, 0.38, 0.92, 12.57, 1.15, NAVY2)
txt(sl, "93 passed", 0.55, 0.98, 5.0, 0.72, size=38, bold=True, color=GREEN)
txt(sl, "0 failed  |  0 errors  |  card-news-pipeline/tests/",
    0.55, 1.65, 12.2, 0.35, size=11, color=RGBColor(0xAA,0xBB,0xDD))

test_modules = [
    ("test_pipeline.py",       "전체 파이프라인 통합 흐름"),
    ("test_selector.py",       "소재 선택 알고리즘 (날짜 우선순위·중복 제거)"),
    ("test_db_source.py",      "DB 로드 / Mock 폴백 동작"),
    ("test_llm_scripter.py",   "Claude Haiku 스크립팅 / 템플릿 폴백"),
    ("test_scripter.py",       "규칙 기반 스크립트 생성"),
    ("test_renderer.py",       "Playwright 슬라이드 렌더링"),
    ("test_orchestrator.py",   "전체 오케스트레이터 흐름"),
    ("test_discord_review.py", "Discord Webhook 검수 요청"),
    ("test_youtube_upload.py", "YouTube OAuth2 업로드"),
    ("test_video.py",          "FFmpeg MP4 조립"),
    ("test_alerts.py",         "오류 알림 전송"),
    ("test_run_log.py",        "실행 이력 로깅"),
    ("test_store.py",          "used_ids 영속화"),
    ("test_review_status.py",  "검수 상태 JSON"),
    ("test_history.py",        "실행 히스토리"),
    ("test_run.py",            "진입점 run.py 흐름"),
    ("test_health_check.py",   "헬스체크 엔드포인트"),
    ("test_mock_data.py",      "샘플 데이터 로드"),
]

per_col = 9
for i, (mod, desc) in enumerate(test_modules):
    ci = i // per_col
    ri = i % per_col
    x = 0.38 + ci * 6.48
    y = 2.22 + ri * 0.54
    bg = LGRAY if ri % 2 == 0 else WHITE
    rect(sl, x, y, 6.3, 0.52, bg)
    txt(sl, mod,  x + 0.12, y + 0.08, 3.1,  0.38, size=9,   color=NAVY, bold=True)
    txt(sl, desc, x + 3.25, y + 0.08, 2.92, 0.38, size=8.5, color=DGRAY)


# ════════════════════════════════════════════════════════════════
# Slide 16 — 결론 · 배운 점 · 개선 방향
# ════════════════════════════════════════════════════════════════
sl = slide()
page_header(sl, "14  결론 · 배운 점 · 개선 방향")
page_footer(sl, 16)

# 결론
rect(sl, 0.38, 0.92, 12.57, 0.42, NAVY)
txt(sl, "결론  —  무엇을 만들었나",
    0.55, 0.97, 12.2, 0.34, size=11, bold=True, color=WHITE)
conclusions = [
    "FastAPI + SQLAlchemy async 풀스택으로 실제 배포 가능한 B2B SaaS를 구현했습니다.",
    "3계층 AI 파이프라인(L1 규칙→L2 Gemini→L3 Claude)으로 비용을 최소화하면서 고위협 탐지 정확도를 높였습니다.",
    "카드뉴스 파이프라인을 통해 탐지된 위협 데이터를 YouTube Shorts 영상까지 완전 자동화했습니다.",
    "Vercel + Railway 듀얼 배포로 서버리스 API와 상시 Celery 워커를 분리 운영하는 구조를 완성했습니다.",
]
for ri, c in enumerate(conclusions):
    bg = LGRAY if ri % 2 == 0 else WHITE
    rect(sl, 0.38, 1.34 + ri * 0.54, 12.57, 0.52, bg)
    txt(sl, "✓  " + c, 0.52, 1.39 + ri * 0.54, 12.3, 0.44, size=9.5, color=NAVY, wrap=True)

# 배운 점
rect(sl, 0.38, 3.58, 5.9, 0.42, BLUE)
txt(sl, "배운 점",
    0.55, 3.63, 5.6, 0.34, size=11, bold=True, color=WHITE)
lessons = [
    "서버리스 환경의 제약 (lifespan·logging·datetime)은\n실제 배포 전에 반드시 검증해야 한다.",
    "AI API 비용은 레이어드 아키텍처로 통제해야\n프로덕션에서 비용 폭발을 막을 수 있다.",
    "외부 API(Polar·Svix·Gemini)의 스펙은\n공식 문서와 실제 동작이 다를 수 있다.",
    "Mock과 실동작을 명확히 구분하지 않으면\n완료 기준이 모호해져 품질이 낮아진다.",
]
for ri, l in enumerate(lessons):
    bg = LGRAY if ri % 2 == 0 else WHITE
    rect(sl, 0.38, 4.0 + ri * 0.75, 5.9, 0.73, bg)
    txt(sl, l, 0.52, 4.06 + ri * 0.75, 5.65, 0.62, size=9.5, color=DGRAY, wrap=True)

# 개선 방향
rect(sl, 6.65, 3.58, 6.3, 0.42, AMBER)
txt(sl, "향후 개선 방향",
    6.82, 3.63, 6.0, 0.34, size=11, bold=True, color=WHITE)
improvements = [
    ("단기 (v1.0)",  "Instagram · TikTok 수집기 구현\n이미지 pHash 분석 파이프라인 통합"),
    ("중기 (v1.5)",  "한국 커뮤니티 실시간 크롤러\nHyperCLOVA X L2 기본 엔진 연동"),
    ("장기 (v2.0)",  "멀티모달 이미지·영상 분석\n화이트라벨 대량 공급 인프라"),
    ("품질 개선",    "오탐률 < 10% 달성 검증\n테스트 커버리지 90%+ 목표"),
]
for ri, (stage, desc) in enumerate(improvements):
    bg = LGRAY if ri % 2 == 0 else WHITE
    rect(sl, 6.65, 4.0 + ri * 0.75, 6.3, 0.73, bg)
    txt(sl, stage, 6.82, 4.06 + ri * 0.75, 1.6,  0.62, size=9, bold=True, color=AMBER)
    txt(sl, desc,  8.45, 4.06 + ri * 0.75, 4.38, 0.62, size=9.5, color=DGRAY, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 17 — 마무리 / Q&A
# ════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, 13.33, 7.5, NAVY)
rect(sl, 0, 0, 13.33, 0.06, BLUE)
rect(sl, 0, 7.44, 13.33, 0.06, BLUE)
rect(sl, 0, 0, 0.18, 7.5, BLUE)

txt(sl, "Q & A",
    0.55, 1.3, 11, 1.6, size=72, bold=True, color=WHITE)
txt(sl, "감사합니다",
    0.55, 3.05, 11, 0.75, size=22,
    color=RGBColor(0xAA, 0xBB, 0xDD))

rect(sl, 0.55, 3.95, 10.8, 0.05, BLUE)

summary = [
    ("프로젝트",  "SAYbrand — AI 기반 브랜드 위협 모니터링 SaaS"),
    ("핵심 기술", "3계층 AI 파이프라인  |  리스크 스코어링  |  카드뉴스 자동 생성"),
    ("배포",      "Vercel (API + Frontend)  +  Railway (Celery Worker)"),
    ("테스트",    "카드뉴스 파이프라인 93 tests passed"),
    ("버전",      "v0.3.1"),
]
for i, (label, val) in enumerate(summary):
    rect(sl, 0.55, 4.15 + i * 0.6, 10.8, 0.58, NAVY2)
    txt(sl, label, 0.72, 4.2 + i * 0.6, 2.0, 0.45, size=10,
        color=RGBColor(0xAA,0xBB,0xDD))
    txt(sl, val,   2.75, 4.2 + i * 0.6, 8.5, 0.45, size=10,
        bold=True, color=WHITE)

txt(sl, TODAY, 0.55, 7.08, 12.5, 0.3, size=8.5,
    color=RGBColor(0x44,0x55,0x77), align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════════
out_path = Path(__file__).parent / "SAYbrand_deck.pptx"
prs.save(str(out_path))
print(f"저장 완료: {out_path}")
print(f"슬라이드 수: {len(prs.slides)}장")
