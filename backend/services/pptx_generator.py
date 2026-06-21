"""PPT 보고서 생성 서비스 — 위협 요약 프레젠테이션"""
from __future__ import annotations

import io
import logging
from sqlalchemy.ext.asyncio import AsyncSession

from backend.services.report_generator import generate_report

logger = logging.getLogger(__name__)


async def generate_pptx_report(
    user_id: int,
    period: str,
    db: AsyncSession,
    org_id: int | None = None,
) -> bytes:
    data = await generate_report(user_id, period, db, org_id=org_id)

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        NAVY  = RGBColor(0x0c, 0x14, 0x28)
        BLUE  = RGBColor(0x1a, 0x6e, 0xf8)
        RED   = RGBColor(0xE2, 0x4B, 0x4A)
        AMBER = RGBColor(0xBA, 0x75, 0x17)
        GREEN = RGBColor(0x1D, 0x9E, 0x75)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LGRAY = RGBColor(0xee, 0xee, 0xee)
        GRAY  = RGBColor(0x88, 0x88, 0x88)
        LBLUE = RGBColor(0xe8, 0xf0, 0xfb)

        s = data["summary"]
        brand      = data.get("brand_info", {}).get("name", "내 브랜드")
        period_range = data["period"]
        gen_date   = data["generated_at"][:10]
        total      = s["total_threats"] or 1

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]  # blank

        PLAT_KO = {
            "youtube": "YouTube", "instagram": "Instagram",
            "x": "X(Twitter)", "tiktok": "TikTok", "naver": "Naver",
        }
        SEV_KO = {
            "critical": "위험", "high": "높음", "medium": "중간", "low": "낮음",
        }
        SEV_COLOR = {
            "critical": RED, "high": AMBER, "medium": BLUE, "low": GREEN,
        }

        def _add_slide():
            return prs.slides.add_slide(blank)

        def _rect(slide, l, t, w, h, fill, no_line=True):
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
            if no_line:
                shp.line.fill.background()
            return shp

        def _txt(slide, text, l, t, w, h, size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = str(text or "")
            run.font.size = Pt(size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color

        def _footer(slide):
            _txt(slide, f"SAYbrand  |  {gen_date}  |  기밀 문서",
                 0.3, 7.15, 12.7, 0.3, size=8, color=GRAY)

        # ────────────────────────────────────────────────────────────
        # Slide 1 — Cover
        # ────────────────────────────────────────────────────────────
        sl1 = _add_slide()
        _rect(sl1, 0, 0, 13.33, 7.5, NAVY)
        _rect(sl1, 0, 0, 13.33, 0.08, BLUE)

        period_type_ko = {"daily": "일간", "weekly": "주간", "monthly": "월간"}.get(
            data.get("period_type", ""), "기간"
        )
        _txt(sl1, "SAYbrand  —  브랜드 위협 모니터링",
             0.8, 0.5, 11, 0.5, size=13, color=RGBColor(0xAA, 0xBB, 0xDD))
        _txt(sl1, brand, 0.8, 1.2, 11, 1.3, size=40, bold=True, color=WHITE)
        _txt(sl1, f"브랜드 위협 {period_type_ko} 보고서",
             0.8, 2.7, 10, 0.7, size=22, color=WHITE)
        _txt(sl1, f"보고 기간 : {period_range}",
             0.8, 3.5, 9, 0.45, size=13, color=RGBColor(0xAA, 0xBB, 0xDD))
        _txt(sl1, f"생성일 : {gen_date}",
             0.8, 4.0, 9, 0.45, size=13, color=RGBColor(0xAA, 0xBB, 0xDD))

        kpi_boxes = [
            ("총 위협", f"{s['total_threats']:,}건"),
            ("미해결", f"{s['unresolved_count']:,}건"),
            ("브랜드 점수", f"{s['brand_score']}/100"),
            ("부정 언급", f"{s['negative_mentions']:,}건"),
        ]
        for i, (lbl, val) in enumerate(kpi_boxes):
            x = 0.8 + i * 3.13
            _rect(sl1, x, 5.1, 2.9, 1.6, RGBColor(0x1d, 0x2a, 0x45))
            _txt(sl1, lbl, x + 0.15, 5.2, 2.6, 0.45, size=10,
                 color=RGBColor(0xAA, 0xBB, 0xDD))
            _txt(sl1, val, x + 0.15, 5.6, 2.6, 0.8, size=24, bold=True, color=WHITE)

        _footer(sl1)

        # ────────────────────────────────────────────────────────────
        # Slide 2 — Executive Summary
        # ────────────────────────────────────────────────────────────
        sl2 = _add_slide()
        _rect(sl2, 0, 0, 13.33, 0.72, NAVY)
        _txt(sl2, "1. 핵심 요약 (Executive Summary)",
             0.4, 0.1, 10, 0.52, size=18, bold=True, color=WHITE)
        _txt(sl2, period_range, 10.0, 0.18, 3.1, 0.4, size=10,
             color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.RIGHT)

        score_val   = s["brand_score"]
        score_color = GREEN if score_val >= 70 else (AMBER if score_val >= 40 else RED)

        kpi6 = [
            ("총 위협 건수",    f"{s['total_threats']:,}건",           NAVY),
            ("미해결 위협",     f"{s['unresolved_count']:,}건",        AMBER),
            ("실제 위협 해결",  f"{s['real_resolved_count']:,}건",     GREEN),
            ("부정적 언급",     f"{s['negative_mentions']:,}건",       RED),
            ("조직적 공격",     f"{s.get('organized_count',0):,}건",   RED),
            ("봇 의심 계정",    f"{s.get('bot_count',0):,}건",         AMBER),
        ]
        for i, (lbl, val, col) in enumerate(kpi6):
            row = i // 3
            c   = i % 3
            x = 0.5 + c * 4.12
            y = 0.9 + row * 2.45
            _rect(sl2, x, y, 3.9, 2.1, RGBColor(0xf1, 0xf5, 0xf9))
            _txt(sl2, lbl, x + 0.15, y + 0.15, 3.6, 0.5, size=11, color=GRAY)
            _txt(sl2, val, x + 0.15, y + 0.6,  3.6, 0.9, size=26, bold=True, color=col)

        _rect(sl2, 0.5, 5.85, 7.1, 1.15, LBLUE)
        _txt(sl2, "브랜드 이미지 점수",
             0.7, 5.95, 5, 0.4, size=11, color=GRAY)
        _txt(sl2, f"{score_val} / 100",
             0.7, 6.25, 5, 0.65, size=22, bold=True, color=score_color)

        _footer(sl2)

        # ────────────────────────────────────────────────────────────
        # Slide 3 — 위협 현황 분석 (심각도 / 플랫폼)
        # ────────────────────────────────────────────────────────────
        sl3 = _add_slide()
        _rect(sl3, 0, 0, 13.33, 0.72, NAVY)
        _txt(sl3, "2. 위협 현황 분석",
             0.4, 0.1, 10, 0.52, size=18, bold=True, color=WHITE)

        by_s = data.get("by_severity", {})
        _txt(sl3, "심각도별 분포", 0.5, 0.85, 5.5, 0.45, size=13, bold=True, color=NAVY)
        y_pos = 1.35
        for sev in ("critical", "high", "medium", "low"):
            cnt = by_s.get(sev, 0)
            if not cnt:
                continue
            ratio = cnt / total
            col = SEV_COLOR.get(sev, GRAY)
            _txt(sl3, f"{SEV_KO.get(sev, sev)}  {cnt:,}건  ({ratio*100:.1f}%)",
                 0.5, y_pos, 3.8, 0.38, size=10, bold=True, color=col)
            bar_w = max(0.05, ratio * 5.5)
            _rect(sl3, 0.5, y_pos + 0.38, bar_w, 0.22, col)
            if ratio < 1.0:
                _rect(sl3, 0.5 + bar_w, y_pos + 0.38, 5.5 - bar_w, 0.22, LGRAY)
            y_pos += 0.78

        by_p = data.get("by_platform", {})
        _txt(sl3, "플랫폼별 분포", 7.0, 0.85, 5.5, 0.45, size=13, bold=True, color=NAVY)
        y_pos2 = 1.35
        for plat, cnt in sorted(by_p.items(), key=lambda x: x[1], reverse=True):
            ratio = cnt / total
            _txt(sl3, f"{PLAT_KO.get(plat, plat)}  {cnt:,}건  ({ratio*100:.1f}%)",
                 7.0, y_pos2, 3.8, 0.38, size=10, color=NAVY)
            bar_w = max(0.05, ratio * 5.5)
            _rect(sl3, 7.0, y_pos2 + 0.38, bar_w, 0.22, BLUE)
            if ratio < 1.0:
                _rect(sl3, 7.0 + bar_w, y_pos2 + 0.38, 5.5 - bar_w, 0.22, LGRAY)
            y_pos2 += 0.78

        _footer(sl3)

        # ────────────────────────────────────────────────────────────
        # Slide 4 — 감성 분석
        # ────────────────────────────────────────────────────────────
        sl4 = _add_slide()
        _rect(sl4, 0, 0, 13.33, 0.72, NAVY)
        _txt(sl4, "3. 감성 · 감정 분석",
             0.4, 0.1, 10, 0.52, size=18, bold=True, color=WHITE)

        by_senti = data.get("by_sentiment", {})
        senti_total = sum(by_senti.values()) or 1
        senti_color_map = {"negative": RED, "neutral": GRAY, "positive": GREEN}
        senti_ko = {"negative": "부정", "neutral": "중립", "positive": "긍정"}

        _txt(sl4, "감성 분포", 0.5, 0.85, 5.5, 0.45, size=13, bold=True, color=NAVY)
        y_pos = 1.35
        for sk in ("negative", "neutral", "positive"):
            cnt = by_senti.get(sk, 0)
            ratio = cnt / senti_total
            col = senti_color_map.get(sk, GRAY)
            _txt(sl4, f"{senti_ko[sk]}  {cnt:,}건  ({ratio*100:.1f}%)",
                 0.5, y_pos, 4, 0.38, size=10, bold=True, color=col)
            bar_w = max(0.05, ratio * 5.5)
            _rect(sl4, 0.5, y_pos + 0.38, bar_w, 0.22, col)
            if ratio < 1.0:
                _rect(sl4, 0.5 + bar_w, y_pos + 0.38, 5.5 - bar_w, 0.22, LGRAY)
            y_pos += 0.78

        by_emo = data.get("by_emotion", {})
        if by_emo:
            emo_total = sum(by_emo.values()) or 1
            _txt(sl4, "감정 분류", 7.0, 0.85, 5.5, 0.45, size=13, bold=True, color=NAVY)
            y_pos2 = 1.35
            for emo, cnt in sorted(by_emo.items(), key=lambda x: x[1], reverse=True)[:6]:
                ratio = cnt / emo_total
                _txt(sl4, f"{emo}  {cnt:,}건  ({ratio*100:.1f}%)",
                     7.0, y_pos2, 4, 0.38, size=10, color=NAVY)
                bar_w = max(0.05, ratio * 5.5)
                _rect(sl4, 7.0, y_pos2 + 0.38, bar_w, 0.22, BLUE)
                if ratio < 1.0:
                    _rect(sl4, 7.0 + bar_w, y_pos2 + 0.38, 5.5 - bar_w, 0.22, LGRAY)
                y_pos2 += 0.78

        _footer(sl4)

        # ────────────────────────────────────────────────────────────
        # Slide 5 — 미해결 위협 Top 5
        # ────────────────────────────────────────────────────────────
        sl5 = _add_slide()
        _rect(sl5, 0, 0, 13.33, 0.72, NAVY)
        _txt(sl5, "4. 미해결 위협 Top 5 (위험도 순)",
             0.4, 0.1, 10, 0.52, size=18, bold=True, color=WHITE)

        unresolved = data.get("unresolved_threats", [])[:5]

        col_x = [0.3, 0.8, 2.1, 3.3, 4.3, 6.2]
        col_w = [0.5, 1.3, 1.2, 1.0, 1.9, 6.5]
        hdrs  = ["#", "플랫폼", "심각도", "위험도", "계정", "내용 미리보기"]

        _rect(sl5, 0.3, 0.82, 12.73, 0.42, NAVY)
        for hdr, x, w in zip(hdrs, col_x, col_w):
            _txt(sl5, hdr, x + 0.05, 0.88, w - 0.1, 0.35, size=9, bold=True, color=WHITE)

        for ri, th in enumerate(unresolved):
            y = 1.25 + ri * 1.15
            bg = RGBColor(0xf8, 0xfa, 0xfc) if ri % 2 == 0 else WHITE
            _rect(sl5, 0.3, y, 12.73, 1.1, bg)

            sev = th.get("severity", "")
            sc  = SEV_COLOR.get(sev, GRAY)
            row_vals = [
                str(ri + 1),
                PLAT_KO.get(th.get("platform", ""), th.get("platform", "")),
                SEV_KO.get(sev, sev),
                str(th.get("risk_score") or 0),
                (th.get("source_account") or "")[:18],
                (th.get("content_preview") or "")[:90],
            ]
            for vi, (val, x, w) in enumerate(zip(row_vals, col_x, col_w)):
                c = sc if vi == 2 else NAVY
                _txt(sl5, val, x + 0.05, y + 0.05, w - 0.1, 1.0,
                     size=9, bold=(vi == 2), color=c)

        if not unresolved:
            _txt(sl5, "해당 기간 미해결 위협이 없습니다.",
                 0.5, 1.5, 12, 0.5, size=13, color=GRAY)

        _footer(sl5)

        # ────────────────────────────────────────────────────────────
        # Slide 6 — 대응 권고사항
        # ────────────────────────────────────────────────────────────
        sl6 = _add_slide()
        _rect(sl6, 0, 0, 13.33, 0.72, NAVY)
        _txt(sl6, "5. 대응 권고사항",
             0.4, 0.1, 10, 0.52, size=18, bold=True, color=WHITE)

        crit_cnt = data.get("by_severity", {}).get("critical", 0)
        high_cnt = data.get("by_severity", {}).get("high", 0)
        org_cnt  = s.get("organized_count", 0)
        bot_cnt  = s.get("bot_count", 0)
        by_senti2 = data.get("by_sentiment", {})
        neg_ratio = by_senti2.get("negative", 0) / (sum(by_senti2.values()) or 1)

        recs = []
        if crit_cnt:
            recs.append((RED, "긴급 대응 필요",
                          f"위험(Critical) {crit_cnt}건 탐지 — 법적 대응 검토 및 공식 입장 발표를 즉시 진행하십시오."))
        if high_cnt:
            recs.append((AMBER, "고위험 위협 모니터링 강화",
                          f"높음(High) {high_cnt}건 — 24시간 집중 모니터링 및 담당자 배정이 필요합니다."))
        if org_cnt:
            recs.append((RED, "조직적 공격 대응",
                          f"조직적 공격 패턴 {org_cnt}건 — SNS 플랫폼 신고 및 법무팀 협의를 병행하십시오."))
        if bot_cnt:
            recs.append((AMBER, "봇 계정 신고",
                          f"봇 의심 계정 {bot_cnt}건 — 각 플랫폼에 신고 및 콘텐츠 삭제 요청을 진행하십시오."))
        if neg_ratio > 0.3:
            recs.append((BLUE, "부정 여론 관리",
                          f"부정 감성 비율 {neg_ratio*100:.1f}% — 고객 소통 강화 및 긍정 콘텐츠 생산을 검토하십시오."))
        if not recs:
            recs.append((GREEN, "정기 모니터링 유지",
                          "현재 심각한 위협이 없습니다. 주기적 모니터링을 유지하십시오."))

        y_pos = 0.85
        for col, title, desc in recs[:5]:
            _rect(sl6, 0.5, y_pos, 12.33, 0.38, col)
            _txt(sl6, f"  {title}", 0.55, y_pos + 0.04, 12, 0.35,
                 size=11, bold=True, color=WHITE)
            _txt(sl6, desc, 0.65, y_pos + 0.44, 12.1, 0.45, size=10, color=NAVY)
            y_pos += 1.05

        _footer(sl6)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    except ImportError:
        logger.warning("python-pptx 미설치 [MOCK] — pip install python-pptx")
        return b""
